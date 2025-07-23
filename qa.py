"""
QA System 

This module provides a document-based question answering system designed for Railway deployment.
It includes functionality for document processing, vector storage, and interaction with OpenRouter's GPT models.
"""

import os
import logging
import threading
import time
import json
import re
import requests
import hashlib
import tempfile
from typing import List, Dict, Any, Optional, Tuple, Union
from datetime import datetime

# Document processing
import pptx
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# ============================================================================
# FIXED DATABASE CONNECTION FOR RAILWAY
# ============================================================================


def get_db_connection() -> Optional[Any]:
    """
    Get database connection for Railway deployment with better error handling.

    Returns:
        Optional[Any]: A database connection object if successful, None otherwise.
    """
    try:
        import psycopg2
        from psycopg2.extras import RealDictCursor

        # Try to get DATABASE_URL first (Railway standard)
        database_url = os.getenv("DATABASE_URL")
        if database_url:
            # Replace postgres:// with postgresql:// if needed
            if database_url.startswith("postgres://"):
                database_url = database_url.replace("postgres://", "postgresql://", 1)
            conn = psycopg2.connect(database_url)
            logger.info("✅ Connected to Railway database via DATABASE_URL")
            return conn

        # Fallback to individual environment variables
        conn = psycopg2.connect(
            host=os.getenv("DB_HOST", "localhost"),
            port=os.getenv("DB_PORT", "5432"),
            database=os.getenv("DB_NAME", "fiftyone_learning"),
            user=os.getenv("DB_USER", "admin"),
            password=os.getenv("DB_PASSWORD", "admin123"),
        )
        logger.info("✅ Connected to database via individual env vars")
        return conn
    except Exception as e:
        logger.error(f"❌ Database connection error: {e}")
        return None


def ensure_stored_documents_table() -> bool:
    """
    Ensure the stored_documents table exists in the database.

    Returns:
        bool: True if the table exists or was created successfully, False otherwise.
    """
    conn = None
    try:
        conn = get_db_connection()
        if not conn:
            return False

        cursor = conn.cursor()

        # Create table if it doesn't exist
        cursor.execute(
            """
            CREATE TABLE IF NOT EXISTS stored_documents (
                id SERIAL PRIMARY KEY,
                filename VARCHAR(255) NOT NULL,
                content BYTEA NOT NULL,
                content_type VARCHAR(100),
                upload_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                UNIQUE(filename)
            )
        """
        )
        conn.commit()
        cursor.close()
        logger.info("✅ Stored documents table ensured")
        return True

    except Exception as e:
        logger.error(f"❌ Error ensuring stored_documents table: {str(e)}")
        return False
    finally:
        if conn:
            conn.close()


def get_documents_from_database() -> List[Tuple[str, bytes, str]]:
    """
    Retrieve documents from the database with proper error handling.

    Returns:
        List[Tuple[str, bytes, str]]: A list of tuples containing (filename, content, content_type)
    """
    conn = None
    try:
        conn = get_db_connection()
        if not conn:
            logger.error("❌ No database connection available")
            return []

        cursor = conn.cursor()

        # Ensure table exists
        if not ensure_stored_documents_table():
            logger.error("❌ Could not ensure stored_documents table exists")
            return []

        # Get documents
        cursor.execute(
            """
            SELECT filename, content, content_type, upload_date
            FROM stored_documents
            ORDER BY upload_date DESC
        """
        )

        documents = cursor.fetchall()
        cursor.close()

        logger.info(f"📄 Retrieved {len(documents)} documents from database")

        # Return in expected format
        result = []
        for filename, content, content_type, upload_date in documents:
            if content:  # Only include documents with content
                result.append((filename, content, content_type))

        return result

    except Exception as e:
        logger.error(f"❌ Error getting documents from database: {str(e)}")
        return []
    finally:
        if conn:
            conn.close()


def create_temp_documents_from_db() -> str:
    """
    Create temporary directory with documents from database.

    Returns:
        str: Path to the temporary directory containing the documents.
    """
    try:
        # Create temporary directory
        temp_dir = tempfile.mkdtemp(prefix="railway_qa_docs_")
        logger.info(f"📁 Created temp directory: {temp_dir}")

        # Get documents from database
        documents = get_documents_from_database()

        if not documents:
            logger.warning("⚠️ No documents found in database")
            # Create a placeholder file so the system doesn't crash
            placeholder_content = """# Welcome to 51Talk AI Learning Platform

This is a placeholder document. Please upload your course materials through the admin panel.

## How to upload documents:
1. Go to /admin/upload_document
2. Upload PDF, PPT, or text files
3. The AI assistant will automatically learn from your materials

## Available features:
- Course material Q&A
- Interactive learning assistance
- Multi-language support
- Progress tracking
"""
            with open(
                os.path.join(temp_dir, "placeholder.txt"), "w", encoding="utf-8"
            ) as f:
                f.write(placeholder_content)
            logger.info("📄 Created placeholder document")
            return temp_dir

        # Save documents to temp directory
        saved_count = 0
        for filename, content, content_type in documents:
            if content:
                try:
                    file_path = os.path.join(temp_dir, filename)
                    with open(file_path, "wb") as f:
                        f.write(content)
                    saved_count += 1
                    logger.info(f"📄 Extracted document: {filename}")
                except Exception as e:
                    logger.error(f"❌ Error saving {filename}: {str(e)}")

        logger.info(
            f"✅ Created temporary documents directory with {saved_count} files"
        )
        return temp_dir

    except Exception as e:
        logger.error(f"❌ Error creating temp documents directory: {str(e)}")
        # Return empty temp directory with placeholder
        temp_dir = tempfile.mkdtemp(prefix="railway_qa_docs_empty_")
        with open(os.path.join(temp_dir, "error.txt"), "w", encoding="utf-8") as f:
            f.write(f"Error loading documents: {str(e)}")
        return temp_dir


def get_openrouter_config() -> Tuple[Optional[str], Optional[str]]:
    """
    Get OpenRouter configuration with better error handling.

    Returns:
        Tuple[Optional[str], Optional[str]]: A tuple containing (api_key, model) or (None, None) if not found.
    """
    try:
        # Try to get from environment variables first
        api_key = os.getenv("OPENROUTER_API_KEY")
        model = os.getenv("OPENROUTER_MODEL", "openai/gpt-4o-mini")

        if api_key:
            logger.info(f"✅ OpenRouter config loaded from env vars - Model: {model}")
            return api_key, model

        # Fallback to config.py
        try:
            from config import get_config

            Config = get_config()

            api_key = getattr(Config, "OPENROUTER_API_KEY", None)
            model = getattr(Config, "OPENROUTER_MODEL", "openai/gpt-4o-mini")

            if api_key:
                logger.info(
                    f"✅ OpenRouter config loaded from config.py - Model: {model}"
                )
                return api_key, model

        except ImportError:
            logger.warning("⚠️ config.py not found, using env vars only")

        logger.error("❌ OPENROUTER_API_KEY not found in environment or config")
        return None, None

    except Exception as e:
        logger.error(f"❌ Error getting OpenRouter config: {e}")
        return None, None


# ============================================================================
# CONVERSATION MEMORY
# ============================================================================


class ConversationMemory:
    """Simple conversation memory for context."""

    def __init__(self, max_history: int = 10) -> None:
        """
        Initialize conversation memory.

        Args:
            max_history (int): Maximum number of messages to retain per user.
        """
        self.max_history = max_history
        self.conversations: Dict[str, List[Dict[str, str]]] = {}

    def add_message(self, user_id: str, question: str, answer: str) -> None:
        """
        Add a message to the conversation history.

        Args:
            user_id (str): Unique identifier for the user.
            question (str): The user's question.
            answer (str): The assistant's response.
        """
        if user_id not in self.conversations:
            self.conversations[user_id] = []
        self.conversations[user_id].append(
            {
                "question": question,
                "answer": answer,
                "timestamp": datetime.now().isoformat(),
            }
        )
        if len(self.conversations[user_id]) > self.max_history:
            self.conversations[user_id] = self.conversations[user_id][
                -self.max_history :
            ]

    def get_context(self, user_id: str, last_n: int = 3) -> str:
        """
        Get recent conversation context for a user.

        Args:
            user_id (str): Unique identifier for the user.
            last_n (int): Number of recent messages to include.

        Returns:
            str: Formatted conversation context.
        """
        if user_id not in self.conversations:
            return ""
        recent = self.conversations[user_id][-last_n:]
        context_parts = []
        for item in recent:
            context_parts.append(f"Human: {item['question']}")
            context_parts.append(f"Assistant: {item['answer']}")
        return "\n".join(context_parts)


# ============================================================================
# OPENROUTER GPT LLM
# ============================================================================


class OpenRouterLLM:
    """OpenRouter GPT API wrapper with improved error handling."""

    def __init__(self, api_key: str, model: str) -> None:
        """
        Initialize the OpenRouter LLM wrapper.

        Args:
            api_key (str): OpenRouter API key.
            model (str): Model identifier to use.
        """
        self.api_key = api_key
        self.model = model
        self.base_url = "https://openrouter.ai/api/v1/chat/completions"
        self.timeout = 60

    def generate_response(self, prompt: str) -> str:
        """
        Generate response with comprehensive error handling.

        Args:
            prompt (str): The prompt to send to the model.

        Returns:
            str: The generated response or an error message.
        """
        if not prompt or not prompt.strip():
            return "[Error] Empty prompt provided"

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://51talk-ai-learning.com",
            "X-Title": "51Talk AI Learning Platform",
        }

        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 800,
            "temperature": 0.7,
            "top_p": 1,
            "frequency_penalty": 0,
            "presence_penalty": 0,
        }

        try:
            response = requests.post(
                self.base_url, headers=headers, json=payload, timeout=self.timeout
            )

            if response.status_code == 429:
                return "[Rate Limited] Too many requests. Please wait a moment and try again."
            elif response.status_code == 401:
                return "[Authentication Error] Invalid OpenRouter API key. Please check your API configuration."
            elif response.status_code == 402:
                return "[Insufficient Credits] Please check your OpenRouter account balance."
            elif response.status_code == 404:
                return f"[Model Error] Model '{self.model}' not found. Please check your model configuration."
            elif response.status_code != 200:
                logger.error(
                    f"OpenRouter API error: {response.status_code} - {response.text}"
                )
                return (
                    f"[API Error] Service unavailable (Status: {response.status_code})"
                )

            data = response.json()

            if "choices" in data and len(data["choices"]) > 0:
                return data["choices"][0]["message"]["content"]
            else:
                return "[Error] No response generated by GPT model"

        except requests.exceptions.Timeout:
            return "[Timeout] GPT request timed out. Please try again."
        except requests.exceptions.ConnectionError:
            return "[Connection Error] Unable to connect to OpenRouter service."
        except Exception as e:
            logger.error(f"OpenRouter API error: {e}")
            return f"[Error] {str(e)}"


# ============================================================================
# DOCUMENT LOADERS
# ============================================================================


class PPTXLoader:
    """Simple PPTX loader."""

    def __init__(self, file_path: str) -> None:
        """
        Initialize the PPTX loader.

        Args:
            file_path (str): Path to the PPTX file.
        """
        self.file_path = file_path

    def load(self) -> List[Dict[str, Any]]:
        """
        Load and extract content from the PPTX file.

        Returns:
            List[Dict[str, Any]]: List of document dictionaries containing slide content and metadata.
        """
        docs = []
        try:
            prs = pptx.Presentation(self.file_path)
            filename = os.path.basename(self.file_path)
            for i, slide in enumerate(prs.slides):
                text_content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_content += cell.text + " "
                            text_content += "\n"
                if text_content.strip():
                    doc = {
                        "page_content": text_content.strip(),
                        "metadata": {
                            "source": self.file_path,
                            "slide_number": i + 1,
                            "page": i + 1,
                            "filename": filename,
                        },
                    }
                    docs.append(doc)
            logger.info(f"📄 Loaded {len(docs)} slides from {filename}")
            return docs
        except Exception as e:
            logger.error(f"❌ Error processing PPTX {self.file_path}: {str(e)}")
            return []


# ============================================================================
# DOCUMENT PROCESSOR - RAILWAY COMPATIBLE
# ============================================================================


class DocumentProcessor:
    """Document processor for Railway deployment."""

    def __init__(self, documents_dir: str = None) -> None:
        """
        Initialize the document processor.

        Args:
            documents_dir (str, optional): Directory containing documents. If None, creates from database.
        """
        self.documents_dir = documents_dir or create_temp_documents_from_db()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800, chunk_overlap=100, length_function=len
        )

    def load_documents(self) -> List[Any]:
        """
        Load documents from the configured directory.

        Returns:
            List[Any]: List of loaded document objects.
        """
        all_docs = []

        # Check if documents directory exists and has files
        if not os.path.exists(self.documents_dir):
            logger.error(f"❌ Documents directory does not exist: {self.documents_dir}")
            return []

        files_found = []
        for root, _, files in os.walk(self.documents_dir):
            for file in files:
                if file.lower().endswith((".pdf", ".pptx", ".ppt", ".txt")):
                    files_found.append(file)

        if not files_found:
            logger.warning(
                f"⚠️ No supported document files found in {self.documents_dir}"
            )
            return []

        logger.info(f"📁 Found {len(files_found)} document files: {files_found}")

        # Process each file
        for root, _, files in os.walk(self.documents_dir):
            for file in files:
                file_path = os.path.join(root, file)
                filename = os.path.basename(file_path)
                try:
                    if file.lower().endswith(".pdf"):
                        docs = self._load_pdf(file_path, filename)
                        all_docs.extend(docs)
                        logger.info(f"📄 Loaded {len(docs)} pages from PDF: {filename}")
                    elif file.lower().endswith((".pptx", ".ppt")):
                        loader = PPTXLoader(file_path)
                        docs = loader.load()
                        all_docs.extend(docs)
                        logger.info(
                            f"📄 Loaded {len(docs)} slides from PPTX: {filename}"
                        )
                    elif file.lower().endswith(".txt"):
                        docs = self._load_txt(file_path, filename)
                        all_docs.extend(docs)
                        logger.info(f"📄 Loaded text file: {filename}")
                except Exception as e:
                    logger.error(f"❌ Error processing {file_path}: {str(e)}")
                    continue

        logger.info(f"✅ Total loaded document sections: {len(all_docs)}")
        return all_docs

    def _load_pdf(self, file_path: str, filename: str) -> List[Any]:
        """
        Load a PDF file.

        Args:
            file_path (str): Path to the PDF file.
            filename (str): Name of the file.

        Returns:
            List[Any]: List of loaded PDF pages.
        """
        try:
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            for doc in docs:
                doc.metadata["filename"] = filename
            return docs
        except Exception as e:
            logger.error(f"❌ Error loading PDF {filename}: {str(e)}")
            return []

    def _load_txt(self, file_path: str, filename: str) -> List[Dict[str, Any]]:
        """
        Load a text file with multiple encoding attempts.

        Args:
            file_path (str): Path to the text file.
            filename (str): Name of the file.

        Returns:
            List[Dict[str, Any]]: List containing a single document dictionary.
        """
        encodings = ["utf-8", "utf-16", "iso-8859-1", "windows-1252"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                if content.strip():
                    return [
                        {
                            "page_content": content.strip(),
                            "metadata": {
                                "source": file_path,
                                "page": 1,
                                "filename": filename,
                            },
                        }
                    ]
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"❌ Error loading TXT {filename}: {str(e)}")
                break
        return []

    def process_documents(self) -> List[Any]:
        """
        Process loaded documents into chunks.

        Returns:
            List[Any]: List of processed document chunks.
        """
        docs = self.load_documents()
        if not docs:
            logger.warning("⚠️ No documents loaded")
            return []

        logger.info(f"📄 Processing {len(docs)} documents into chunks...")
        try:
            texts = []
            metadatas = []
            for doc in docs:
                if hasattr(doc, "page_content"):
                    content = doc.page_content
                    metadata = doc.metadata
                else:
                    content = doc["page_content"]
                    metadata = doc["metadata"]
                content = re.sub(r"\s+", " ", content).strip()
                if len(content) > 50:
                    texts.append(content)
                    metadatas.append(metadata)

            splits = self.text_splitter.create_documents(texts, metadatas=metadatas)
            logger.info(f"✅ Created {len(splits)} document chunks")
            return splits
        except Exception as e:
            logger.error(f"❌ Error processing documents: {str(e)}")
            return []


# ============================================================================
# VECTOR STORE - RAILWAY COMPATIBLE
# ============================================================================


class VectorStore:
    """Vector store manager for Railway deployment."""

    def __init__(self, documents_dir: str = None, vector_db_path: str = "vector_db") -> None:
        """
        Initialize the vector store manager.

        Args:
            documents_dir (str, optional): Directory containing documents. If None, creates from database.
            vector_db_path (str): Path to store vector database.
        """
        self.documents_dir = documents_dir or create_temp_documents_from_db()
        self.vector_db_path = vector_db_path
        self.embeddings = None
        self.vector_store_instance = None
        self.document_chunks = []
        self._initialization_lock = threading.Lock()
        self._is_initializing = False
        self._initialization_complete = False
        self._initialization_error = None
        self._start_background_initialization()

    def _start_background_initialization(self) -> None:
        """Start background initialization of the vector store."""
        def initialize():
            try:
                with self._initialization_lock:
                    if self._initialization_complete or self._is_initializing:
                        return
                    self._is_initializing = True
                    logger.info(
                        "🚀 Starting Railway-compatible vector store initialization..."
                    )

                    # Load documents from temporary directory (created from database)
                    processor = DocumentProcessor(self.documents_dir)
                    self.document_chunks = processor.process_documents()

                    # Mark as complete - we don't need actual vector store for OpenRouter/GPT
                    self._initialization_complete = True
                    self._is_initializing = False
                    logger.info(
                        f"✅ Vector store initialized with {len(self.document_chunks)} chunks"
                    )

            except Exception as e:
                self._initialization_error = str(e)
                self._is_initializing = False
                logger.error(f"❌ Failed to initialize vector store: {str(e)}")

        thread = threading.Thread(target=initialize, daemon=True)
        thread.start()

    def get_vector_store(self, timeout: int = 30) -> "MockRetriever":
        """
        Get a vector store retriever.

        Args:
            timeout (int): Timeout in seconds for initialization.

        Returns:
            MockRetriever: A mock retriever instance.

        Raises:
            TimeoutError: If initialization times out.
            Exception: If initialization fails.
        """
        start_time = time.time()
        while time.time() - start_time < timeout:
            if self._initialization_complete:
                return MockRetriever(self.document_chunks)
            elif self._initialization_error:
                raise Exception(f"Vector store failed: {self._initialization_error}")
            time.sleep(0.5)
        raise TimeoutError("Vector store initialization timed out")

    def is_ready(self) -> bool:
        """
        Check if the vector store is ready.

        Returns:
            bool: True if initialization is complete, False otherwise.
        """
        return self._initialization_complete


class MockRetriever:
    """Mock retriever for Railway deployment."""

    def __init__(self, document_chunks: List[Any]) -> None:
        """
        Initialize the mock retriever.

        Args:
            document_chunks (List[Any]): List of document chunks to search through.
        """
        self.document_chunks = document_chunks
        logger.info(f"🔍 MockRetriever initialized with {len(document_chunks)} chunks")

    def as_retriever(self, search_kwargs: Optional[Dict[str, Any]] = None) -> "MockRetriever":
        """
        Return self as a retriever.

        Args:
            search_kwargs (Optional[Dict[str, Any]]): Optional search parameters.

        Returns:
            MockRetriever: This instance.
        """
        return self

    def get_relevant_documents(
        self, query: str, search_kwargs: Optional[Dict[str, Any]] = None
    ) -> List[Any]:
        """
        Search for relevant documents based on a query.

        Args:
            query (str): The search query.
            search_kwargs (Optional[Dict[str, Any]]): Optional search parameters.

        Returns:
            List[Any]: List of relevant documents.
        """
        if not self.document_chunks:
            logger.warning("⚠️ No document chunks available for search")
            return []

        query_lower = query.lower()
        scored_docs = []

        for doc in self.document_chunks:
            content = (
                doc.page_content
                if hasattr(doc, "page_content")
                else doc.get("page_content", "")
            )
            if not content:
                continue

            content_lower = content.lower()

            # Improved scoring based on keyword matches
            score = 0
            query_words = query_lower.split()
            for word in query_words:
                if len(word) > 2:  # Skip short words
                    # Exact word matches
                    score += content_lower.count(word) * 2
                    # Partial matches
                    if word in content_lower:
                        score += 1

            # Bonus for question words
            question_words = ["what", "how", "why", "when", "where", "which", "who"]
            for qword in question_words:
                if qword in query_lower and qword in content_lower:
                    score += 3

            if score > 0:
                scored_docs.append((score, doc))

        # Sort by score and return top results
        scored_docs.sort(key=lambda x: x[0], reverse=True)
        k = search_kwargs.get("k", 6) if search_kwargs else 6

        results = [doc for score, doc in scored_docs[:k]]
        logger.info(
            f"🔍 Found {len(results)} relevant documents for query: {query[:50]}..."
        )
        return results


# ============================================================================
# MAIN QA SYSTEM - RAILWAY COMPATIBLE
# ============================================================================


class DocumentQA:
    """Main Document QA system for Railway deployment."""

    def __init__(self, documents_dir: str = None) -> None:
        """
        Initialize the Document QA system.

        Args:
            documents_dir (str, optional): Directory containing documents. If None, creates from database.
        """
        self.documents_dir = documents_dir or create_temp_documents_from_db()
        self.vector_store_manager = VectorStore(self.documents_dir)
        self.gpt_llm = None
        self.conversation_memory = ConversationMemory()
        self.response_cache = {}
        self.cache_max_size = 50

        # Initialize OpenRouter GPT
        api_key, model = get_openrouter_config()
        if api_key and model:
            self.gpt_llm = OpenRouterLLM(api_key, model)
            logger.info(
                f"✅ OpenRouter GPT initialized successfully with model: {model}"
            )
        else:
            logger.error("❌ OpenRouter API not configured properly")

        self.greetings = {
            "hi": "Hello! I'm your AI learning assistant. I can help you with course materials, explain concepts, or just chat about your studies. What would you like to know?",
            "hello": "Hi there! I'm here to help with your learning journey. You can ask me about course topics, request explanations, or get study tips. How can I assist you today?",
            "hey": "Hey! Great to see you here. I'm your AI tutor ready to help with anything you're studying. What's on your mind?",
            "good morning": "Good morning! Ready to learn something new today? I'm here to help with your course materials.",
            "good afternoon": "Good afternoon! How's your learning going today? I'm here to help with explanations or questions.",
            "good evening": "Good evening! Perfect time for some learning. What would you like to explore?",
        }

    def answer_question(self, question: str, user_id: Optional[str] = None) -> Dict[str, Any]:
        """
        Answer a question using the QA system.

        Args:
            question (str): The question to answer.
            user_id (Optional[str]): Optional user identifier for conversation history.

        Returns:
            Dict[str, Any]: A dictionary containing:
                - answer: The generated answer
                - sources: List of source documents
                - conversation_type: Type of conversation
        """
        if not question or not question.strip():
            return {
                "answer": "Hello! I'm your AI assistant. I'm here and ready to help! What would you like to know or discuss?",
                "sources": [],
                "conversation_type": "greeting",
            }

        question = question.strip()
        question_lower = question.lower()

        logger.info(f"🤖 Processing question: {question[:50]}...")

        try:
            conversation_type = self._detect_conversation_type(question_lower)
            if conversation_type == "greeting":
                response = self._handle_greeting(question_lower)
                if user_id:
                    self.conversation_memory.add_message(user_id, question, response)
                return {
                    "answer": response,
                    "sources": [],
                    "conversation_type": "greeting",
                }

            response = self._handle_document_question(question, user_id)
            if user_id:
                self.conversation_memory.add_message(
                    user_id, question, response["answer"]
                )
            return response

        except Exception as e:
            logger.error(f"❌ Error in answer_question: {str(e)}")
            import traceback

            logger.error(f"❌ Full traceback: {traceback.format_exc()}")
            return {
                "answer": "I encountered an error, but I'm still here to help! Could you try rephrasing your question?",
                "sources": [],
                "conversation_type": "error",
            }

    def _detect_conversation_type(self, question_lower: str) -> str:
        """
        Detect the type of conversation based on the question.

        Args:
            question_lower (str): The question in lowercase.

        Returns:
            str: Type of conversation ("greeting", "document_based", or "general")
        """
        greeting_patterns = [
            r"\b(hi|hello|hey|hiya)\b",
            r"\bgood (morning|afternoon|evening|night)\b",
            r"\bhow are you\b",
            r"\bwhat\'?s up\b",
        ]
        if any(re.search(pattern, question_lower) for pattern in greeting_patterns):
            return "greeting"

        document_keywords = [
            "explain",
            "definition",
            "what is",
            "how does",
            "tell me about",
            "course material",
            "lesson",
            "chapter",
            "according to",
        ]
        if any(keyword in question_lower for keyword in document_keywords):
            return "document_based"
        return "general"

    def _handle_greeting(self, question_lower: str) -> str:
        """
        Handle greeting questions.

        Args:
            question_lower (str): The greeting question in lowercase.

        Returns:
            str: Appropriate greeting response.
        """
        for greeting, response in self.greetings.items():
            if greeting in question_lower:
                return response
        return "Hello! I'm your AI learning assistant. I'm here to help you understand course materials, answer questions, or discuss topics you're studying. What would you like to explore today?"

    def _handle_document_question(
        self, question: str, user_id: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Handle document-based questions.

        Args:
            question (str): The question to answer.
            user_id (Optional[str]): Optional user identifier.

        Returns:
            Dict[str, Any]: Response dictionary containing answer and metadata.
        """
        try:
            logger.info("🔍 Starting document question processing...")

            # Check 1: OpenRouter GPT availability
            if not self.gpt_llm:
                logger.error("❌ OpenRouter GPT is not configured")
                return {
                    "answer": "❌ OpenRouter GPT is not properly configured. Please check your OPENROUTER_API_KEY environment variable.",
                    "sources": [],
                    "conversation_type": "error",
                }
            logger.info("✅ OpenRouter GPT is available")

            # Check 2: Vector store readiness
            if not self.vector_store_manager.is_ready():
                logger.error("❌ Vector store is not ready")
                if self.vector_store_manager._is_initializing:
                    return {
                        "answer": "⏳ I'm currently loading course materials from the database. Please wait a moment and try again.",
                        "sources": [],
                        "conversation_type": "loading",
                    }
                else:
                    error_msg = self.vector_store_manager._initialization_error
                    logger.error(f"❌ Vector store error: {error_msg}")
                    return {
                        "answer": "📄 I don't have any course materials uploaded yet. Please upload some documents through the admin panel (/admin/upload_document), and I'll be able to help you with course-specific questions!",
                        "sources": [],
                        "conversation_type": "no_documents",
                    }
            logger.info("✅ Vector store is ready")

            # Check 3: Get retriever
            try:
                retriever = self.vector_store_manager.get_vector_store(timeout=10)
                if not retriever:
                    logger.error("❌ Failed to get retriever")
                    return {
                        "answer": "❌ I'm having trouble accessing course materials. Please try again.",
                        "sources": [],
                        "conversation_type": "error",
                    }
                logger.info("✅ Retriever obtained successfully")
            except Exception as retriever_error:
                logger.error(f"❌ Retriever error: {str(retriever_error)}")
                return {
                    "answer": "📄 I don't have any course materials uploaded yet. Please upload some documents through the admin panel (/admin/upload_document), and I'll be able to help you with course-specific questions!",
                    "sources": [],
                    "conversation_type": "no_documents",
                }

            # Check 4: Get relevant documents
            try:
                logger.info("🔍 Searching for relevant documents...")
                docs = retriever.get_relevant_documents(
                    question, search_kwargs={"k": 6}
                )
                logger.info(f"✅ Found {len(docs)} relevant documents")
            except Exception as search_error:
                logger.error(f"❌ Document search error: {str(search_error)}")
                # Fall back to general knowledge
                answer = self._generate_gpt_response_general(question, user_id)
                return {"answer": answer, "sources": [], "conversation_type": "general"}

            # Enhanced search with keywords
            keywords = self._extract_keywords(question)
            if keywords:
                keyword_query = " ".join(keywords)
                try:
                    additional_docs = retriever.get_relevant_documents(
                        keyword_query, search_kwargs={"k": 4}
                    )
                    seen_contents = {
                        (
                            doc.page_content
                            if hasattr(doc, "page_content")
                            else doc.get("page_content", "")
                        )
                        for doc in docs
                    }
                    for doc in additional_docs:
                        content = (
                            doc.page_content
                            if hasattr(doc, "page_content")
                            else doc.get("page_content", "")
                        )
                        if content not in seen_contents:
                            docs.append(doc)
                            seen_contents.add(content)
                except Exception as keyword_error:
                    logger.warning(f"⚠️ Keyword search failed: {str(keyword_error)}")

            # Check 5: Handle no documents case
            if not docs:
                logger.info("⚠️ No relevant documents found")
                documents_in_db = get_documents_from_database()
                if not documents_in_db:
                    return {
                        "answer": "📄 I don't have any course materials uploaded yet. Please upload some documents through the admin panel (/admin/upload_document), and I'll be able to help you with course-specific questions!",
                        "sources": [],
                        "conversation_type": "no_documents",
                    }
                else:
                    logger.info("🤖 Using GPT general knowledge (no document matches)")
                    answer = self._generate_gpt_response_general(question, user_id)
                    return {
                        "answer": answer,
                        "sources": [],
                        "conversation_type": "general",
                    }

            # Check 6: Prepare context and generate response
            try:
                logger.info("📝 Preparing context from documents...")
                context = self._prepare_context(docs)
                sources = self._extract_sources(docs)
                logger.info(
                    f"✅ Context prepared, length: {len(context)} chars, sources: {len(sources)}"
                )

                logger.info("🤖 Generating GPT response...")
                answer = self._generate_gpt_response(question, context, user_id)
                logger.info(f"✅ GPT response generated, length: {len(answer)} chars")

                return {
                    "answer": answer,
                    "sources": sources,
                    "conversation_type": "document_based",
                }
            except Exception as generation_error:
                logger.error(f"❌ Response generation error: {str(generation_error)}")
                # Fall back to general knowledge
                answer = self._generate_gpt_response_general(question, user_id)
                return {"answer": answer, "sources": [], "conversation_type": "general"}

        except Exception as e:
            logger.error(f"❌ Outer exception in document question handling: {str(e)}")
            import traceback

            logger.error(f"❌ Full traceback: {traceback.format_exc()}")
            return {
                "answer": "I encountered an error, but I'm still here to help! Could you try rephrasing your question?",
                "sources": [],
                "conversation_type": "error",
            }

    def _generate_gpt_response_general(self, question: str, user_id: Optional[str] = None) -> str:
        """
        Generate response using GPT without course materials.

        Args:
            question (str): The question to answer.
            user_id (Optional[str]): Optional user identifier.

        Returns:
            str: The generated response.
        """
        if not self.gpt_llm:
            return "[Error: OpenRouter not configured. Please set OPENROUTER_API_KEY environment variable.]"

        prompt = f"""You are a helpful AI learning assistant for the 51Talk AI Learning Platform. 

Please provide a helpful, educational response to this question: {question}

Keep your response clear, educational, and helpful for a student."""

        try:
            response = self.gpt_llm.generate_response(prompt)
            if (
                not response
                or response.startswith("[Error")
                or response.startswith("[Authentication")
            ):
                return "I'm having trouble connecting to my language model. Please check that the OpenRouter API key is configured correctly and has sufficient credits."

            response += "\n\n*Note: This response is based on general knowledge. Upload course materials for more specific help!*"
            return response
        except Exception as e:
            logger.error(f"❌ Error generating general GPT response: {e}")
            return f"I encountered an issue generating a response. Please check the OpenRouter API configuration. Error: {str(e)}"

    def _extract_keywords(self, question: str) -> List[str]:
        """
        Extract keywords from a question.

        Args:
            question (str): The question to analyze.

        Returns:
            List[str]: List of extracted keywords.
        """
        import string

        stop_words = {
            "the",
            "a",
            "an",
            "and",
            "or",
            "but",
            "in",
            "on",
            "at",
            "to",
            "for",
            "of",
            "with",
            "by",
            "is",
            "are",
            "was",
            "were",
            "be",
            "been",
            "being",
            "have",
            "has",
            "had",
            "do",
            "does",
            "did",
            "will",
            "would",
            "could",
            "should",
            "may",
            "might",
            "must",
            "can",
            "what",
            "where",
            "when",
            "why",
            "how",
            "who",
            "which",
            "this",
            "that",
            "these",
            "those",
            "i",
            "you",
            "he",
            "she",
            "it",
            "we",
            "they",
            "me",
            "him",
            "her",
            "us",
            "them",
            "my",
            "your",
            "his",
            "her",
            "its",
            "our",
            "their",
        }
        question_clean = question.lower().translate(
            str.maketrans("", "", string.punctuation)
        )
        words = question_clean.split()
        keywords = [word for word in words if word not in stop_words and len(word) > 2]
        return keywords[:5]

    def _prepare_context(self, docs: List[Any], max_length: int = 3000) -> str:
        """
        Prepare context from document chunks.

        Args:
            docs (List[Any]): List of document chunks.
            max_length (int): Maximum length of context in characters.

        Returns:
            str: Combined context string.
        """
        context_parts = []
        total_length = 0
        for doc in docs:
            content = (
                doc.page_content
                if hasattr(doc, "page_content")
                else doc.get("page_content", "")
            )
            if total_length + len(content) > max_length:
                break
            context_parts.append(content)
            total_length += len(content)
        return "\n\n".join(context_parts)

    def _extract_sources(self, docs: List[Any]) -> List[Dict[str, str]]:
        """
        Extract source information from documents.

        Args:
            docs (List[Any]): List of document chunks.

        Returns:
            List[Dict[str, str]]: List of source dictionaries.
        """
        sources = []
        seen_sources = set()
        for doc in docs:
            metadata = (
                doc.metadata if hasattr(doc, "metadata") else doc.get("metadata", {})
            )
            filename = metadata.get(
                "filename", os.path.basename(metadata.get("source", "Unknown"))
            )
            page = metadata.get("page", metadata.get("slide_number", "Unknown"))
            source_key = f"{filename}_{page}"
            if source_key not in seen_sources:
                sources.append({"file": filename, "page": str(page)})
                seen_sources.add(source_key)
        return sources[:4]

    def _get_cache_key(self, question: str, context: str) -> str:
        """
        Generate a cache key for a question-context pair.

        Args:
            question (str): The question.
            context (str): The context.

        Returns:
            str: MD5 hash of the combined string.
        """
        combined = f"{question}||{context}"
        return hashlib.md5(combined.encode()).hexdigest()

    def _generate_gpt_response(
        self, question: str, context: str, user_id: Optional[str] = None
    ) -> str:
        """
        Generate response using OpenRouter GPT.

        Args:
            question (str): The question to answer.
            context (str): Context from documents.
            user_id (Optional[str]): Optional user identifier.

        Returns:
            str: The generated response.
        """
        if not self.gpt_llm:
            return "[Error: OpenRouter GPT not configured. Please check your API settings.]"

        # Check cache
        cache_key = self._get_cache_key(question, context)
        if cache_key in self.response_cache:
            logger.debug("💾 Cache hit for question")
            return self.response_cache[cache_key]

        # Create prompt optimized for GPT models
        prompt = f"""You are a helpful AI learning assistant for the 51Talk AI Learning Platform. You provide clear, accurate, and educational responses based on course materials.

Guidelines:
- Use the course materials provided to give accurate, helpful answers
- Be conversational and friendly in your tone
- Provide specific examples from the course materials when relevant
- If the materials don't fully answer the question, acknowledge this clearly
- Keep responses educational and informative, tailored to the learning context
- Focus on helping students understand concepts rather than just providing facts

Course Materials:
{context}

Student Question: {question}

Please provide a helpful, educational response based on the course materials above:"""

        try:
            response = self.gpt_llm.generate_response(prompt)

            if (
                not response
                or response.startswith("[Error")
                or response.startswith("[Authentication")
            ):
                return "I'm having trouble connecting to my language model. Please check that the OpenRouter API key is configured correctly and has sufficient credits."

            # Cache response
            if len(self.response_cache) >= self.cache_max_size:
                # Remove oldest entry
                oldest_key = next(iter(self.response_cache))
                del self.response_cache[oldest_key]
            self.response_cache[cache_key] = response

            return response
        except Exception as e:
            logger.error(f"❌ Error generating GPT response: {e}")
            return "[Error] I encountered an issue generating a response. Please check your OpenRouter API configuration."

    def get_status(self) -> Dict[str, Any]:
        """
        Get system status information.

        Returns:
            Dict[str, Any]: Dictionary containing system status information.
        """
        # Count documents in database
        documents = get_documents_from_database()
        document_count = len(documents)

        return {
            "ready": self.vector_store_manager.is_ready() and self.gpt_llm is not None,
            "initializing": self.vector_store_manager._is_initializing,
            "error": self.vector_store_manager._initialization_error,
            "document_count": document_count,
            "llm_provider": "OpenRouter GPT",
            "gpt_available": self.gpt_llm is not None,
            "cache_size": len(self.response_cache),
            "deployment": "Railway Compatible",
            "documents_dir": self.documents_dir,
            "chunks_loaded": len(self.vector_store_manager.document_chunks),
        }


# ============================================================================
# GLOBAL INSTANCE MANAGEMENT
# ============================================================================

_qa_instance = None


def initialize_qa(
    documents_dir: Optional[str] = None, llama_model_path: Optional[str] = None
) -> DocumentQA:
    """
    Initialize the QA system for Railway deployment.

    Args:
        documents_dir (Optional[str]): Optional directory containing documents.
        llama_model_path (Optional[str]): Optional path to LLaMA model (unused in this implementation).

    Returns:
        DocumentQA: The initialized QA system instance.
    """
    global _qa_instance
    if _qa_instance is None:
        # For Railway, always create temp directory from database
        temp_dir = documents_dir or create_temp_documents_from_db()
        logger.info(f"🚀 Initializing Railway-compatible QA system from: {temp_dir}")
        _qa_instance = DocumentQA(temp_dir)
    return _qa_instance


def get_qa_system() -> Optional[DocumentQA]:
    """
    Get the global QA instance.

    Returns:
        Optional[DocumentQA]: The QA system instance if initialized, None otherwise.
    """
    return _qa_instance


def ask_question(
    question: str,
    documents_dir: Optional[str] = None,
    llama_model_path: Optional[str] = None,
    user_id: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Simple function to ask a question using GPT on Railway.

    Args:
        question (str): The question to ask.
        documents_dir (Optional[str]): Optional directory containing documents.
        llama_model_path (Optional[str]): Optional path to LLaMA model (unused).
        user_id (Optional[str]): Optional user identifier.

    Returns:
        Dict[str, Any]: Response dictionary containing answer and metadata.
    """
    qa_system = initialize_qa(documents_dir)
    return qa_system.answer_question(question, user_id)


def get_system_status(
    documents_dir: Optional[str] = None, llama_model_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    Get system status for Railway deployment.

    Args:
        documents_dir (Optional[str]): Optional directory containing documents.
        llama_model_path (Optional[str]): Optional path to LLaMA model (unused).

    Returns:
        Dict[str, Any]: Dictionary containing system status information.
    """
    qa_system = initialize_qa(documents_dir)
    return qa_system.get_status()


def cleanup_qa() -> None:
    """Cleanup function for graceful shutdown."""
    global _qa_instance
    try:
        if _qa_instance:
            logger.info("🧹 Cleaning up Railway QA system resources...")
            _qa_instance.response_cache.clear()
            _qa_instance.conversation_memory.conversations.clear()
            _qa_instance = None
            logger.info("✅ Railway QA cleanup completed")
    except Exception as e:
        logger.error(f"❌ Error during QA cleanup: {str(e)}")


# ============================================================================
# MAIN EXECUTION
# ============================================================================

if __name__ == "__main__":
    try:
        print("🚀 Starting QA System Test with OpenRouter GPT...")
        qa = DocumentQA()
        status = qa.get_status()
        print("QA System Status:", json.dumps(status, indent=2))

        if status["ready"]:
            response = qa.answer_question("Hello!")
            print(f"\n👋 Greeting Response: {response['answer']}")

            response = qa.answer_question("What is artificial intelligence?")
            print(f"\n🤖 AI Question Response: {response['answer']}")
        else:
            print("⏳ QA system is not ready. Check your configuration.")

    except Exception as e:
        print(f"❌ Test failed: {str(e)}")
        import traceback

        traceback.print_exc()